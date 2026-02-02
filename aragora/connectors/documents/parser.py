"""
Unified Document Parser for Omnivorous Document Ingestion.

Provides consistent parsing interface for various document formats:
- PDF (with table extraction)
- Word documents (.docx)
- Excel spreadsheets (.xlsx, .xls)
- PowerPoint presentations (.pptx)
- Plain text, Markdown, HTML
- Structured data (JSON, YAML, XML, CSV)

Usage:
    from aragora.connectors.documents import DocumentParser, parse_document

    # Auto-detect format from extension
    result = await parse_document(content, filename="report.pdf")

    # Or use parser directly
    parser = DocumentParser()
    result = parser.parse(content, format="pdf")
"""

from __future__ import annotations

import io
import json
import logging
import mimetypes
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import IO, Any, Optional, cast

logger = logging.getLogger(__name__)


class DocumentFormat(str, Enum):
    """Supported document formats."""

    # Office documents
    PDF = "pdf"
    DOCX = "docx"
    DOC = "doc"
    XLSX = "xlsx"
    XLS = "xls"
    PPTX = "pptx"
    PPT = "ppt"

    # Text formats
    TXT = "txt"
    MD = "md"
    RST = "rst"
    HTML = "html"
    RTF = "rtf"

    # Structured data
    JSON = "json"
    YAML = "yaml"
    XML = "xml"
    CSV = "csv"

    # Code
    CODE = "code"

    # Notebooks and interactive formats
    IPYNB = "ipynb"  # Jupyter notebooks

    # E-books
    EPUB = "epub"  # Electronic publication
    MOBI = "mobi"  # Kindle format

    # Archives
    ZIP = "zip"
    TAR = "tar"
    GZIP = "gzip"

    UNKNOWN = "unknown"


@dataclass
class DocumentChunk:
    """A chunk of extracted content from a document."""

    content: str
    chunk_type: str = "text"  # text, table, heading, code, etc.
    page: int | None = None
    section: str | None = None
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass
class DocumentTable:
    """A table extracted from a document."""

    data: list[list[str]]  # 2D array of cell values
    headers: Optional[list[str]] = None  # Column headers if detected
    page: int | None = None  # Page number where table was found
    caption: str | None = None  # Table caption if available
    metadata: dict[str, Any] = field(default_factory=dict)

    @property
    def rows(self) -> int:
        """Number of rows in the table."""
        return len(self.data)

    @property
    def cols(self) -> int:
        """Number of columns in the table."""
        return len(self.data[0]) if self.data else 0

    def to_markdown(self) -> str:
        """Convert table to markdown format."""
        if not self.data:
            return ""

        lines = []
        headers = self.headers or self.data[0]
        lines.append("| " + " | ".join(str(h) for h in headers) + " |")
        lines.append("| " + " | ".join("---" for _ in headers) + " |")

        start_row = 0 if self.headers else 1
        for row in self.data[start_row:]:
            lines.append("| " + " | ".join(str(cell) for cell in row) + " |")

        return "\n".join(lines)


@dataclass
class ParsedDocument:
    """Result of parsing a document."""

    content: str  # Full text content
    chunks: list[DocumentChunk] = field(default_factory=list)
    format: DocumentFormat = DocumentFormat.UNKNOWN
    filename: str | None = None  # Original filename
    title: str | None = None  # Document title if extracted
    metadata: dict[str, Any] = field(default_factory=dict)
    tables: list[DocumentTable | list[list[str]]] = field(default_factory=list)
    pages: int = 0
    word_count: int = 0
    errors: list[str] = field(default_factory=list)

    def __post_init__(self):
        if not self.word_count and self.content:
            self.word_count = len(self.content.split())


# Extension to format mapping
EXTENSION_MAP: dict[str, DocumentFormat] = {
    ".pdf": DocumentFormat.PDF,
    ".docx": DocumentFormat.DOCX,
    ".doc": DocumentFormat.DOC,
    ".xlsx": DocumentFormat.XLSX,
    ".xls": DocumentFormat.XLS,
    ".pptx": DocumentFormat.PPTX,
    ".ppt": DocumentFormat.PPT,
    ".txt": DocumentFormat.TXT,
    ".md": DocumentFormat.MD,
    ".markdown": DocumentFormat.MD,
    ".rst": DocumentFormat.RST,
    ".html": DocumentFormat.HTML,
    ".htm": DocumentFormat.HTML,
    ".rtf": DocumentFormat.RTF,
    ".json": DocumentFormat.JSON,
    ".yaml": DocumentFormat.YAML,
    ".yml": DocumentFormat.YAML,
    ".xml": DocumentFormat.XML,
    ".csv": DocumentFormat.CSV,
    ".py": DocumentFormat.CODE,
    ".js": DocumentFormat.CODE,
    ".ts": DocumentFormat.CODE,
    ".java": DocumentFormat.CODE,
    ".go": DocumentFormat.CODE,
    ".rs": DocumentFormat.CODE,
    ".rb": DocumentFormat.CODE,
    ".php": DocumentFormat.CODE,
    # Notebooks
    ".ipynb": DocumentFormat.IPYNB,
    # E-books
    ".epub": DocumentFormat.EPUB,
    ".mobi": DocumentFormat.MOBI,
    # Archives
    ".zip": DocumentFormat.ZIP,
    ".tar": DocumentFormat.TAR,
    ".gz": DocumentFormat.GZIP,
    ".tgz": DocumentFormat.TAR,
    ".tar.gz": DocumentFormat.TAR,
}


class DocumentParser:
    """
    Unified document parser for omnivorous ingestion.

    Supports multiple document formats with consistent output.
    """

    def __init__(
        self,
        max_pages: int = 100,
        max_content_size: int = 500_000,  # 500KB text limit
        extract_tables: bool = True,
        extract_metadata: bool = True,
    ):
        """
        Initialize document parser.

        Args:
            max_pages: Maximum pages to extract from PDFs
            max_content_size: Maximum content size in bytes
            extract_tables: Whether to extract tables separately
            extract_metadata: Whether to extract document metadata
        """
        self.max_pages = max_pages
        self.max_content_size = max_content_size
        self.extract_tables = extract_tables
        self.extract_metadata = extract_metadata

    def detect_format(
        self, filename: str | None = None, content: bytes | None = None
    ) -> DocumentFormat:
        """Detect document format from filename or content."""
        if filename:
            ext = Path(filename).suffix.lower()
            if ext in EXTENSION_MAP:
                return EXTENSION_MAP[ext]

        # Try MIME type detection
        if filename:
            mime_type, _ = mimetypes.guess_type(filename)
            if mime_type:
                if "pdf" in mime_type:
                    return DocumentFormat.PDF
                if "word" in mime_type or "msword" in mime_type:
                    return DocumentFormat.DOCX
                if "spreadsheet" in mime_type or "excel" in mime_type:
                    return DocumentFormat.XLSX
                if "presentation" in mime_type or "powerpoint" in mime_type:
                    return DocumentFormat.PPTX

        # Magic number detection for binary content
        if content:
            if content[:4] == b"%PDF":
                return DocumentFormat.PDF
            if content[:4] == b"PK\x03\x04":  # ZIP-based (Office 2007+)
                # Could be docx, xlsx, pptx - check further if needed
                return DocumentFormat.DOCX  # Default to docx
            if content[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":  # OLE
                return DocumentFormat.DOC  # Legacy Office

        return DocumentFormat.UNKNOWN

    def parse(
        self,
        content: bytes | str,
        format: DocumentFormat | None = None,
        filename: str | None = None,
    ) -> ParsedDocument:
        """
        Parse a document and extract content.

        Args:
            content: Document content (bytes or string)
            format: Document format (auto-detected if not provided)
            filename: Original filename (for format detection)

        Returns:
            ParsedDocument with extracted content
        """
        # Convert string to bytes if needed
        if isinstance(content, str):
            content_bytes = content.encode("utf-8")
        else:
            content_bytes = content

        # Detect format
        if format is None:
            format = self.detect_format(filename, content_bytes)

        # Route to appropriate parser
        try:
            if format == DocumentFormat.PDF:
                return self._parse_pdf(content_bytes)
            elif format in (DocumentFormat.DOCX, DocumentFormat.DOC):
                return self._parse_docx(content_bytes)
            elif format in (DocumentFormat.XLSX, DocumentFormat.XLS):
                return self._parse_excel(content_bytes)
            elif format in (DocumentFormat.PPTX, DocumentFormat.PPT):
                return self._parse_pptx(content_bytes)
            elif format == DocumentFormat.HTML:
                return self._parse_html(content_bytes)
            elif format == DocumentFormat.JSON:
                return self._parse_json(content_bytes)
            elif format == DocumentFormat.YAML:
                return self._parse_yaml(content_bytes)
            elif format == DocumentFormat.XML:
                return self._parse_xml(content_bytes)
            elif format == DocumentFormat.CSV:
                return self._parse_csv(content_bytes)
            elif format == DocumentFormat.IPYNB:
                return self._parse_jupyter(content_bytes)
            elif format == DocumentFormat.EPUB:
                return self._parse_epub(content_bytes)
            elif format == DocumentFormat.MOBI:
                return self._parse_mobi(content_bytes)
            elif format in (DocumentFormat.ZIP, DocumentFormat.TAR, DocumentFormat.GZIP):
                return self._parse_archive(content_bytes, format)
            else:
                # Default to plain text
                return self._parse_text(content_bytes, format)
        except (ValueError, KeyError, IndexError, TypeError, UnicodeDecodeError, OSError) as e:
            logger.error(f"Failed to parse document: {e}")
            return ParsedDocument(
                content="",
                format=format,
                errors=[str(e)],
            )

    def _parse_pdf(self, content: bytes) -> ParsedDocument:
        """Parse PDF document."""
        chunks: list[DocumentChunk] = []
        tables: list[DocumentTable | list[list[str]]] = []
        all_text: list[str] = []
        metadata: dict[str, Any] = {}
        errors: list[str] = []
        pages = 0

        # Try pdfplumber first (better table extraction)
        try:
            import pdfplumber

            with pdfplumber.open(io.BytesIO(content)) as pdf:
                pages = len(pdf.pages)
                if self.extract_metadata and pdf.metadata:
                    metadata = dict(pdf.metadata)

                for i, page in enumerate(pdf.pages[: self.max_pages]):
                    # Extract text
                    text = page.extract_text() or ""
                    if text:
                        all_text.append(text)
                        chunks.append(
                            DocumentChunk(
                                content=text,
                                chunk_type="page",
                                page=i + 1,
                            )
                        )

                    # Extract tables
                    if self.extract_tables:
                        page_tables = page.extract_tables()
                        for table in page_tables:
                            if table:
                                tables.append(table)
                                # Also add table as chunk
                                table_text = "\n".join(
                                    "\t".join(str(cell) if cell else "" for cell in row)
                                    for row in table
                                )
                                chunks.append(
                                    DocumentChunk(
                                        content=table_text,
                                        chunk_type="table",
                                        page=i + 1,
                                    )
                                )

            return ParsedDocument(
                content="\n\n".join(all_text)[: self.max_content_size],
                chunks=chunks,
                format=DocumentFormat.PDF,
                metadata=metadata,
                tables=tables,
                pages=pages,
                errors=errors,
            )

        except ImportError:
            pass  # Fall through to pypdf
        except (ValueError, KeyError, OSError, TypeError) as e:
            errors.append(f"pdfplumber error: {e}")

        # Fallback to pypdf
        try:
            try:
                from pypdf import PdfReader
            except ImportError:
                from PyPDF2 import PdfReader

            reader = PdfReader(io.BytesIO(content))
            pages = len(reader.pages)

            if self.extract_metadata and reader.metadata:
                metadata = {
                    k.lstrip("/"): v for k, v in reader.metadata.items() if isinstance(v, str)
                }

            for i, page in enumerate(reader.pages[: self.max_pages]):
                text = page.extract_text() or ""
                if text:
                    all_text.append(text)
                    chunks.append(
                        DocumentChunk(
                            content=text,
                            chunk_type="page",
                            page=i + 1,
                        )
                    )

            return ParsedDocument(
                content="\n\n".join(all_text)[: self.max_content_size],
                chunks=chunks,
                format=DocumentFormat.PDF,
                metadata=metadata,
                pages=pages,
                errors=errors,
            )

        except ImportError:
            errors.append("No PDF parser available (install pypdf or pdfplumber)")
        except (ValueError, KeyError, OSError, TypeError) as e:
            errors.append(f"pypdf error: {e}")

        return ParsedDocument(
            content="",
            format=DocumentFormat.PDF,
            errors=errors,
        )

    def _parse_docx(self, content: bytes) -> ParsedDocument:
        """Parse Word document (.docx)."""
        chunks: list[DocumentChunk] = []
        tables: list[DocumentTable | list[list[str]]] = []
        all_text: list[str] = []
        metadata: dict[str, Any] = {}
        errors: list[str] = []

        try:
            from docx import Document

            doc = Document(io.BytesIO(content))

            # Extract metadata
            if self.extract_metadata and doc.core_properties:
                props = doc.core_properties
                metadata = {
                    "author": props.author,
                    "title": props.title,
                    "subject": props.subject,
                    "created": str(props.created) if props.created else None,
                    "modified": str(props.modified) if props.modified else None,
                }
                metadata = {k: v for k, v in metadata.items() if v}

            # Extract paragraphs
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    all_text.append(text)
                    # Detect headings
                    chunk_type = (
                        "heading" if para.style and "Heading" in para.style.name else "paragraph"
                    )
                    chunks.append(
                        DocumentChunk(
                            content=text,
                            chunk_type=chunk_type,
                            metadata={"style": para.style.name if para.style else None},
                        )
                    )

            # Extract tables
            if self.extract_tables:
                for table in doc.tables:
                    table_data: list[list[str]] = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    if table_data:
                        tables.append(table_data)
                        table_text = "\n".join("\t".join(row) for row in table_data)
                        chunks.append(DocumentChunk(content=table_text, chunk_type="table"))

            return ParsedDocument(
                content="\n\n".join(all_text)[: self.max_content_size],
                chunks=chunks,
                format=DocumentFormat.DOCX,
                metadata=metadata,
                tables=tables,
                pages=len(doc.sections),
                errors=errors,
            )

        except ImportError:
            errors.append("python-docx not installed (pip install python-docx)")
        except (ValueError, KeyError, OSError, TypeError, AttributeError) as e:
            errors.append(f"docx parse error: {e}")

        return ParsedDocument(
            content="",
            format=DocumentFormat.DOCX,
            errors=errors,
        )

    def _parse_excel(self, content: bytes) -> ParsedDocument:
        """Parse Excel spreadsheet (.xlsx, .xls)."""
        chunks: list[DocumentChunk] = []
        tables: list[DocumentTable | list[list[str]]] = []
        all_text: list[str] = []
        metadata: dict[str, Any] = {}
        errors: list[str] = []

        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)

            # Extract metadata
            if self.extract_metadata and wb.properties:
                props = wb.properties
                metadata = {
                    "creator": props.creator,
                    "title": props.title,
                    "subject": props.subject,
                    "created": str(props.created) if props.created else None,
                    "modified": str(props.modified) if props.modified else None,
                    "sheets": wb.sheetnames,
                }
                metadata = {k: v for k, v in metadata.items() if v}

            # Extract content from each sheet
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_data: list[list[str]] = []

                for row in sheet.iter_rows(max_row=1000):  # Limit rows
                    row_data = [str(cell.value) if cell.value is not None else "" for cell in row]
                    if any(row_data):  # Skip empty rows
                        sheet_data.append(row_data)

                if sheet_data:
                    tables.append(sheet_data)
                    sheet_text = f"Sheet: {sheet_name}\n" + "\n".join(
                        "\t".join(row) for row in sheet_data
                    )
                    all_text.append(sheet_text)
                    chunks.append(
                        DocumentChunk(
                            content=sheet_text,
                            chunk_type="sheet",
                            section=sheet_name,
                            metadata={"sheet": sheet_name, "rows": len(sheet_data)},
                        )
                    )

            wb.close()

            return ParsedDocument(
                content="\n\n".join(all_text)[: self.max_content_size],
                chunks=chunks,
                format=DocumentFormat.XLSX,
                metadata=metadata,
                tables=tables,
                pages=len(wb.sheetnames),
                errors=errors,
            )

        except ImportError:
            errors.append("openpyxl not installed (pip install openpyxl)")
        except (ValueError, KeyError, OSError, TypeError, IndexError) as e:
            errors.append(f"Excel parse error: {e}")

        return ParsedDocument(
            content="",
            format=DocumentFormat.XLSX,
            errors=errors,
        )

    def _parse_pptx(self, content: bytes) -> ParsedDocument:
        """Parse PowerPoint presentation (.pptx)."""
        chunks: list[DocumentChunk] = []
        all_text: list[str] = []
        metadata: dict[str, Any] = {}
        errors: list[str] = []

        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))

            # Extract metadata
            if self.extract_metadata and prs.core_properties:
                props = prs.core_properties
                metadata = {
                    "author": props.author,
                    "title": props.title,
                    "subject": props.subject,
                    "created": str(props.created) if props.created else None,
                    "modified": str(props.modified) if props.modified else None,
                }
                metadata = {k: v for k, v in metadata.items() if v}

            # Extract content from slides
            for i, slide in enumerate(prs.slides):
                slide_texts: list[str] = []
                notes_text = ""

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text)

                    # Extract table content
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text for cell in row.cells]
                            slide_texts.append("\t".join(row_text))

                # Extract speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text

                slide_content = "\n".join(slide_texts)
                if slide_content:
                    all_text.append(f"Slide {i + 1}:\n{slide_content}")
                    chunks.append(
                        DocumentChunk(
                            content=slide_content,
                            chunk_type="slide",
                            page=i + 1,
                            metadata={"notes": notes_text} if notes_text else {},
                        )
                    )

                if notes_text:
                    chunks.append(
                        DocumentChunk(
                            content=notes_text,
                            chunk_type="notes",
                            page=i + 1,
                        )
                    )

            return ParsedDocument(
                content="\n\n".join(all_text)[: self.max_content_size],
                chunks=chunks,
                format=DocumentFormat.PPTX,
                metadata=metadata,
                pages=len(prs.slides),
                errors=errors,
            )

        except ImportError:
            errors.append("python-pptx not installed (pip install python-pptx)")
        except (ValueError, KeyError, OSError, TypeError, AttributeError) as e:
            errors.append(f"PowerPoint parse error: {e}")

        return ParsedDocument(
            content="",
            format=DocumentFormat.PPTX,
            errors=errors,
        )

    def _parse_html(self, content: bytes) -> ParsedDocument:
        """Parse HTML content."""
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(content, "html.parser")

            # Remove script and style elements
            for element in soup(["script", "style", "nav", "footer", "header"]):
                element.decompose()

            text = soup.get_text(separator="\n", strip=True)

            # Extract metadata
            metadata: dict[str, Any] = {}
            if self.extract_metadata:
                title = soup.find("title")
                if title:
                    metadata["title"] = title.get_text()
                for meta in soup.find_all("meta"):
                    name = meta.get("name") or meta.get("property")
                    content_val = meta.get("content")
                    if name and content_val:
                        metadata[name] = content_val

            return ParsedDocument(
                content=text[: self.max_content_size],
                format=DocumentFormat.HTML,
                metadata=metadata,
            )

        except ImportError:
            return ParsedDocument(
                content=content.decode("utf-8", errors="ignore")[: self.max_content_size],
                format=DocumentFormat.HTML,
                errors=["beautifulsoup4 not installed for HTML parsing"],
            )
        except (ValueError, TypeError, UnicodeDecodeError) as e:
            return ParsedDocument(
                content="",
                format=DocumentFormat.HTML,
                errors=[str(e)],
            )

    def _parse_json(self, content: bytes) -> ParsedDocument:
        """Parse JSON content."""
        try:
            text = content.decode("utf-8")
            data = json.loads(text)

            # Pretty print for readability
            formatted = json.dumps(data, indent=2, ensure_ascii=False)

            return ParsedDocument(
                content=formatted[: self.max_content_size],
                format=DocumentFormat.JSON,
                metadata={"keys": list(data.keys()) if isinstance(data, dict) else []},
            )
        except (json.JSONDecodeError, UnicodeDecodeError, TypeError) as e:
            return ParsedDocument(
                content=content.decode("utf-8", errors="ignore")[: self.max_content_size],
                format=DocumentFormat.JSON,
                errors=[str(e)],
            )

    def _parse_yaml(self, content: bytes) -> ParsedDocument:
        """Parse YAML content."""
        try:
            import yaml

            text = content.decode("utf-8")
            data = yaml.safe_load(text)

            # Convert back to YAML for consistent formatting
            formatted = yaml.dump(data, default_flow_style=False, allow_unicode=True)

            return ParsedDocument(
                content=formatted[: self.max_content_size],
                format=DocumentFormat.YAML,
                metadata={"keys": list(data.keys()) if isinstance(data, dict) else []},
            )
        except ImportError:
            return ParsedDocument(
                content=content.decode("utf-8", errors="ignore")[: self.max_content_size],
                format=DocumentFormat.YAML,
                errors=["pyyaml not installed"],
            )
        except (ValueError, TypeError, UnicodeDecodeError) as e:
            return ParsedDocument(
                content=content.decode("utf-8", errors="ignore")[: self.max_content_size],
                format=DocumentFormat.YAML,
                errors=[str(e)],
            )

    def _parse_xml(self, content: bytes) -> ParsedDocument:
        """Parse XML content safely (defusedxml prevents XXE attacks)."""
        try:
            import defusedxml.ElementTree as ET

            root = ET.fromstring(content)

            # Extract all text content
            texts: list[str] = []

            def extract_text(element, depth=0):
                if element.text and element.text.strip():
                    texts.append(element.text.strip())
                for child in element:
                    extract_text(child, depth + 1)
                if element.tail and element.tail.strip():
                    texts.append(element.tail.strip())

            extract_text(root)

            return ParsedDocument(
                content="\n".join(texts)[: self.max_content_size],
                format=DocumentFormat.XML,
                metadata={"root_tag": root.tag},
            )
        except (ValueError, TypeError, UnicodeDecodeError) as e:
            return ParsedDocument(
                content=content.decode("utf-8", errors="ignore")[: self.max_content_size],
                format=DocumentFormat.XML,
                errors=[str(e)],
            )

    def _parse_csv(self, content: bytes) -> ParsedDocument:
        """Parse CSV content."""
        try:
            import csv

            text = content.decode("utf-8")
            reader = csv.reader(io.StringIO(text))
            rows = list(reader)

            # Format as tab-separated for readability
            formatted = "\n".join("\t".join(row) for row in rows)

            return ParsedDocument(
                content=formatted[: self.max_content_size],
                format=DocumentFormat.CSV,
                tables=[rows],
                metadata={"rows": len(rows), "columns": len(rows[0]) if rows else 0},
            )
        except (ValueError, TypeError, UnicodeDecodeError, IndexError) as e:
            return ParsedDocument(
                content=content.decode("utf-8", errors="ignore")[: self.max_content_size],
                format=DocumentFormat.CSV,
                errors=[str(e)],
            )

    def _parse_jupyter(self, content: bytes) -> ParsedDocument:
        """Parse Jupyter notebook (.ipynb) files."""
        chunks: list[DocumentChunk] = []
        all_text: list[str] = []
        metadata: dict[str, Any] = {}
        errors: list[str] = []

        try:
            notebook = json.loads(content.decode("utf-8"))
            metadata = notebook.get("metadata", {})

            # Extract kernel info
            kernel_info = metadata.get("kernelspec", {})
            if kernel_info:
                metadata["kernel"] = kernel_info.get("display_name", "Unknown")

            cells = notebook.get("cells", [])
            for i, cell in enumerate(cells):
                cell_type = cell.get("cell_type", "code")
                source = cell.get("source", [])

                # Handle source as list or string
                if isinstance(source, list):
                    cell_content = "".join(source)
                else:
                    cell_content = source

                if not cell_content.strip():
                    continue

                # Format based on cell type
                if cell_type == "markdown":
                    all_text.append(cell_content)
                    chunks.append(
                        DocumentChunk(
                            content=cell_content,
                            chunk_type="markdown",
                            metadata={"cell_index": i, "cell_type": "markdown"},
                        )
                    )
                elif cell_type == "code":
                    # Include code with syntax marker
                    code_block = f"```python\n{cell_content}\n```"
                    all_text.append(code_block)
                    chunks.append(
                        DocumentChunk(
                            content=cell_content,
                            chunk_type="code",
                            metadata={"cell_index": i, "cell_type": "code"},
                        )
                    )

                    # Also extract outputs if present
                    outputs = cell.get("outputs", [])
                    for output in outputs:
                        output_type = output.get("output_type", "")
                        if output_type == "stream":
                            text = "".join(output.get("text", []))
                            if text.strip():
                                all_text.append(f"Output:\n{text}")
                        elif output_type in ("execute_result", "display_data"):
                            data = output.get("data", {})
                            if "text/plain" in data:
                                text = "".join(data["text/plain"])
                                if text.strip():
                                    all_text.append(f"Output:\n{text}")
                elif cell_type == "raw":
                    all_text.append(cell_content)
                    chunks.append(
                        DocumentChunk(
                            content=cell_content,
                            chunk_type="text",
                            metadata={"cell_index": i, "cell_type": "raw"},
                        )
                    )

            return ParsedDocument(
                content="\n\n".join(all_text)[: self.max_content_size],
                chunks=chunks,
                format=DocumentFormat.IPYNB,
                metadata=metadata,
                errors=errors,
            )

        except (json.JSONDecodeError, KeyError, TypeError, UnicodeDecodeError) as e:
            logger.error(f"Failed to parse Jupyter notebook: {e}")
            return ParsedDocument(
                content="",
                format=DocumentFormat.IPYNB,
                errors=[str(e)],
            )

    def _parse_epub(self, content: bytes) -> ParsedDocument:
        """Parse EPUB e-book files."""
        chunks: list[DocumentChunk] = []
        all_text: list[str] = []
        metadata: dict[str, Any] = {}
        errors: list[str] = []

        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup

            book = epub.read_epub(io.BytesIO(content))

            # Extract metadata
            if book.get_metadata("DC", "title"):
                metadata["title"] = book.get_metadata("DC", "title")[0][0]
            if book.get_metadata("DC", "creator"):
                metadata["author"] = book.get_metadata("DC", "creator")[0][0]
            if book.get_metadata("DC", "language"):
                metadata["language"] = book.get_metadata("DC", "language")[0][0]
            if book.get_metadata("DC", "publisher"):
                metadata["publisher"] = book.get_metadata("DC", "publisher")[0][0]

            # Extract chapters/content
            chapter_num = 0
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    chapter_num += 1
                    html_content = item.get_content().decode("utf-8", errors="ignore")
                    soup = BeautifulSoup(html_content, "html.parser")

                    # Get chapter title if available
                    title = soup.find(["h1", "h2", "h3"])
                    chapter_title = title.get_text().strip() if title else f"Chapter {chapter_num}"

                    # Extract text
                    text = soup.get_text(separator="\n", strip=True)
                    if text:
                        all_text.append(f"## {chapter_title}\n\n{text}")
                        chunks.append(
                            DocumentChunk(
                                content=text,
                                chunk_type="chapter",
                                section=chapter_title,
                                metadata={"chapter": chapter_num},
                            )
                        )

            return ParsedDocument(
                content="\n\n".join(all_text)[: self.max_content_size],
                chunks=chunks,
                format=DocumentFormat.EPUB,
                title=metadata.get("title"),
                metadata=metadata,
                pages=chapter_num,
                errors=errors,
            )

        except ImportError:
            errors.append(
                "ebooklib not installed. Install with: pip install ebooklib beautifulsoup4"
            )
            return ParsedDocument(
                content="",
                format=DocumentFormat.EPUB,
                errors=errors,
            )
        except (ValueError, KeyError, OSError, TypeError, UnicodeDecodeError) as e:
            logger.error(f"Failed to parse EPUB: {e}")
            return ParsedDocument(
                content="",
                format=DocumentFormat.EPUB,
                errors=[str(e)],
            )

    def _parse_mobi(self, content: bytes) -> ParsedDocument:
        """Parse MOBI/Kindle e-book files."""
        errors: list[str] = []

        try:
            import mobi

            # mobi library extracts to a temp directory
            import tempfile
            import os

            with tempfile.TemporaryDirectory() as tmpdir:
                mobi_path = os.path.join(tmpdir, "book.mobi")
                with open(mobi_path, "wb") as f:
                    f.write(content)

                # Extract MOBI to get HTML content
                tempdir, extracted = mobi.extract(mobi_path)

                # Read the extracted HTML (with path traversal protection)
                html_path = os.path.realpath(os.path.join(tempdir, extracted))
                if not html_path.startswith(os.path.realpath(tempdir)):
                    return ParsedDocument(
                        content="",
                        format=DocumentFormat.MOBI,
                        errors=["MOBI extraction produced invalid path"],
                    )
                if os.path.exists(html_path):
                    with open(html_path, "rb") as f:
                        html_content = f.read()

                    # Parse as HTML
                    return self._parse_html(html_content)

            return ParsedDocument(
                content="",
                format=DocumentFormat.MOBI,
                errors=["Could not extract MOBI content"],
            )

        except ImportError:
            errors.append("mobi not installed. Install with: pip install mobi")
            return ParsedDocument(
                content="",
                format=DocumentFormat.MOBI,
                errors=errors,
            )
        except (ValueError, KeyError, OSError, TypeError, UnicodeDecodeError) as e:
            logger.error(f"Failed to parse MOBI: {e}")
            return ParsedDocument(
                content="",
                format=DocumentFormat.MOBI,
                errors=[str(e)],
            )

    def _parse_archive(self, content: bytes, format: DocumentFormat) -> ParsedDocument:
        """Parse archive files (ZIP, TAR, GZIP) and extract documents."""
        chunks: list[DocumentChunk] = []
        all_text: list[str] = []
        metadata: dict[str, Any] = {"files": []}
        errors: list[str] = []
        tables: list[DocumentTable | list[list[str]]] = []

        try:
            import zipfile
            import tarfile
            import gzip

            extracted_docs: list[ParsedDocument] = []

            if format == DocumentFormat.ZIP:
                with zipfile.ZipFile(io.BytesIO(content), "r") as zf:
                    for name in zf.namelist():
                        if name.endswith("/"):  # Skip directories
                            continue
                        # Check if it's a supported format
                        ext = Path(name).suffix.lower()
                        if ext in EXTENSION_MAP and ext not in (".zip", ".tar", ".gz"):
                            try:
                                file_content = zf.read(name)
                                doc = self.parse(file_content, filename=name)
                                doc.filename = name  # Ensure filename is set
                                extracted_docs.append(doc)
                                metadata["files"].append(name)
                            except (OSError, ValueError, KeyError, zipfile.BadZipFile) as e:
                                errors.append(f"Failed to extract {name}: {e}")

            elif format in (DocumentFormat.TAR, DocumentFormat.GZIP):
                # Handle .tar.gz and .tar
                mode: str = (
                    "r:gz" if format == DocumentFormat.GZIP or content[:2] == b"\x1f\x8b" else "r"
                )
                try:
                    with tarfile.open(  # type: ignore[call-overload]
                        fileobj=cast(IO[bytes], io.BytesIO(content)),
                        mode=mode,  # type: ignore[arg-type]  # tarfile Literal mode types are stricter than str
                    ) as tf:
                        for member in tf.getmembers():
                            if member.isfile():
                                ext = Path(member.name).suffix.lower()
                                if ext in EXTENSION_MAP and ext not in (
                                    ".zip",
                                    ".tar",
                                    ".gz",
                                    ".tgz",
                                ):
                                    try:
                                        f = tf.extractfile(member)
                                        if f:
                                            file_content = f.read()
                                            doc = self.parse(file_content, filename=member.name)
                                            doc.filename = member.name  # Ensure filename is set
                                            extracted_docs.append(doc)
                                            metadata["files"].append(member.name)
                                    except (OSError, ValueError, KeyError, tarfile.TarError) as e:
                                        errors.append(f"Failed to extract {member.name}: {e}")
                except tarfile.ReadError:
                    # Try as plain gzip (single file)
                    try:
                        decompressed = gzip.decompress(content)
                        # Try to parse as text
                        return self._parse_text(decompressed)
                    except (OSError, ValueError, gzip.BadGzipFile) as e:
                        errors.append(f"Failed to decompress gzip: {e}")

            # Combine all extracted documents
            for doc in extracted_docs:
                if doc.content:
                    all_text.append(f"--- File: {doc.filename or 'unknown'} ---\n{doc.content}")
                chunks.extend(doc.chunks)
                tables.extend(doc.tables)
                if doc.errors:
                    errors.extend(doc.errors)

            metadata["file_count"] = len(extracted_docs)

            return ParsedDocument(
                content="\n\n".join(all_text)[: self.max_content_size],
                chunks=chunks,
                format=format,
                metadata=metadata,
                tables=tables,
                errors=errors,
            )

        except (OSError, ValueError, KeyError, zipfile.BadZipFile, tarfile.TarError) as e:
            logger.error(f"Failed to parse archive: {e}")
            return ParsedDocument(
                content="",
                format=format,
                errors=[str(e)],
            )

    def _parse_text(
        self, content: bytes, format: DocumentFormat = DocumentFormat.TXT
    ) -> ParsedDocument:
        """Parse plain text content."""
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("latin-1", errors="ignore")

        return ParsedDocument(
            content=text[: self.max_content_size],
            format=format,
        )


# Convenience function
async def parse_document(
    content: bytes | str | Path,
    filename: str | None = None,
    format: DocumentFormat | None = None,
    **kwargs,
) -> ParsedDocument:
    """
    Parse a document asynchronously.

    Args:
        content: Document content (bytes, string, or path)
        filename: Original filename for format detection
        format: Document format (auto-detected if not provided)
        **kwargs: Additional options for DocumentParser

    Returns:
        ParsedDocument with extracted content
    """
    # Read from path if needed
    if isinstance(content, Path):
        filename = filename or content.name
        content = content.read_bytes()

    parser = DocumentParser(**kwargs)
    return parser.parse(content, format=format, filename=filename)


__all__ = [
    "DocumentParser",
    "DocumentFormat",
    "DocumentChunk",
    "DocumentTable",
    "ParsedDocument",
    "parse_document",
    "EXTENSION_MAP",
]
